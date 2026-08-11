# -*- coding: utf-8 -*-
"""Gera OTC_Tracker_One_Pager.pptx — dois slides 16:9: o que a solução faz e o
mapa de cobertura dos processos (ferramenta a ferramenta, produto a produto).

Uso:

    pip install python-pptx           # uma vez (não é dependência do app)
    python scripts/build_onepager_pptx.py
    python scripts/build_onepager_pptx.py /tmp/previa.html   # + prévia HTML

Paleta e geometria são as da própria aplicação (`visual-refresh.css`): gradiente
#0066cc → #5e5ce6 → #8b5cf6 → #d946ef, cartão branco com borda leve e canto
arredondado — o slide tem de parecer o sistema.

O conteúdo é CAPACIDADE e FLUXO, mais os benefícios já entregues (a pedido:
consolidação das Intelligent Solutions e a redução de 0,5 FTE, no rodapé).
Para editar o texto, mexa em `PILARES` e `COLUNAS` e rode de novo — o arquivo é
gerado, não editado à mão, senão a próxima rodada apaga a edição.
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)), '..'))
OUT = os.path.join(ROOT, 'OTC_Tracker_One_Pager.pptx')

# ── Paleta (visual-refresh.css) ──────────────────────────────────────────────
A1, A2, A3, A4 = '0066CC', '5E5CE6', '8B5CF6', 'D946EF'
INK = RGBColor(0x0E, 0x11, 0x2A)
INK_SOFT = RGBColor(0x54, 0x5A, 0x72)
INK_FAINT = RGBColor(0x7B, 0x82, 0x99)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BD = RGBColor(0xE3, 0xE6, 0xEF)
PAGE_BG = RGBColor(0xF6, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Segoe UI'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])          # em branco


def _no_line(shape):
    shape.line.fill.background()


def rect(x, y, w, h, fill=None, radius=None, line=None, line_w=Pt(1)):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if radius is not None:
        # adj = raio / menor lado
        shp.adjustments[0] = radius / min(w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        _no_line(shp)
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    shp.text_frame.word_wrap = True
    return shp


def gradient(shape, stops, angle_deg=100):
    """Gradiente linear multi-parada — python-pptx só expõe duas, então o
    preenchimento é escrito direto no XML da forma."""
    spPr = shape._element.spPr
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    gs = ''.join(
        '<a:gs pos="%d"><a:srgbClr val="%s"/></a:gs>' % (int(pos * 1000), color)
        for pos, color in stops)
    xml = ('<a:gradFill %s rotWithShape="1"><a:gsLst>%s</a:gsLst>'
           '<a:lin ang="%d" scaled="0"/></a:gradFill>'
           % (nsdecls('a'), gs, int(angle_deg * 60000)))
    frag = parse_xml(xml)
    # o preenchimento vem logo depois da geometria
    ref = spPr.find(qn('a:prstGeom'))
    ref.addnext(frag)


def text(shape, blocks, margins=(0.14, 0.10, 0.14, 0.10), anchor=MSO_ANCHOR.TOP):
    """blocks: [(texto, tamanho, negrito, cor, espaço_antes, alinhamento, espaçamento)]"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    l, t, r, b = margins
    tf.margin_left, tf.margin_top = Inches(l), Inches(t)
    tf.margin_right, tf.margin_bottom = Inches(r), Inches(b)
    for i, blk in enumerate(blocks):
        txt, size, bold, color, space, align, spacing = (list(blk) + [None] * 7)[:7]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align or PP_ALIGN.LEFT
        if space:
            p.space_before = Pt(space)
        p.line_spacing = spacing or 1.0
        run = p.add_run()
        run.text = txt
        f = run.font
        f.name, f.size, f.bold = FONT, Pt(size), bool(bold)
        f.color.rgb = color
    return shape


def textbox(x, y, w, h, blocks, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return text(box, blocks, margins=(0, 0, 0, 0), anchor=anchor)


# ── Fundo ────────────────────────────────────────────────────────────────────
rect(0, 0, 13.333, 7.5, fill=PAGE_BG)

# ── Faixa do cabeçalho ───────────────────────────────────────────────────────
band = rect(0, 0, 13.333, 1.42)
gradient(band, [(0, A1), (0.5, A2), (0.8, A3), (1, A4)])

textbox(0.62, 0.24, 8.6, 0.95, [
    ('OTC Tracker', 30, True, WHITE, None, None, 1.0),
    ('One platform for the full OTC derivatives lifecycle — Brazil OTC Operations',
     14, False, RGBColor(0xE4, 0xEC, 0xFF), 3, None, 1.0),
])
textbox(7.25, 0.56, 5.45, 0.4, [
    ('Trade capture  ›  B3 / CETIP registration  ›  Client confirmation  ›  Settlement',
     9.5, True, RGBColor(0xEB, 0xF1, 0xFF), None, PP_ALIGN.RIGHT, 1.0),
])

# ── Quatro pilares ───────────────────────────────────────────────────────────
PILARES = [
    (A1, 'One solution, end to end',
     'Registration, confirmation, settlement and regulatory reporting live in the '
     'same place — every user reads and writes the same data.',
     ['100+ screens covering NDF, FX and commodity options, swaps, unwinds and intragroup legs',
      'The step that follows always knows what the previous one did — no re-keying between tools',
      'No side spreadsheet as the source of truth — nothing left to reconcile between tools']),
    (A2, 'The routine runs itself',
     'What used to be typed, copied or remembered is now imported, generated or '
     'scheduled — people are left with the decisions.',
     ['Trades pulled from Athena automatically — NDF every 20 minutes, FX options hourly',
      'Client confirmations, settlement advices and B3 / CETIP files produced by the solution',
      'Daily routines e-mail each user exactly what is still open, without anyone asking']),
    (A3, 'Controls built into the flow',
     'The check is part of the screen, not a separate review — so nothing depends '
     'on someone remembering to do it.',
     ['Confirmation trail OTC › MO / FO, with each stage signed only by its own users',
      'Deadlines in business days; signing off late requires a written justification',
      'Every action stamped with date, time and employee ID, plus maker / checker on the data']),
    (A4, 'Changes without a release',
     'The rules the users own are edited on screen and take effect on the next '
     'click — no code change, no deployment window.',
     ['25 on-screen registries hold the mappings that used to be hard-coded',
      'Access granted page by page — and routine by routine inside the Control Panel',
      'One code base serving users in English, Portuguese and Spanish']),
]

CARD_Y, CARD_H, CARD_W, GAP, X0 = 1.62, 3.50, 2.98, 0.24, 0.62
for i, (accent, titulo, lead, bullets) in enumerate(PILARES):
    x = X0 + i * (CARD_W + GAP)
    card = rect(x, CARD_Y, CARD_W, CARD_H, fill=CARD_BG, radius=0.17, line=CARD_BD, line_w=Pt(0.75))
    # faixa de acento no topo do cartão
    tab = rect(x + 0.22, CARD_Y + 0.28, 0.52, 0.075, fill=RGBColor.from_string(accent), radius=0.037)
    textbox(x + 0.22, CARD_Y + 0.48, CARD_W - 0.44, 0.62, [
        (titulo, 15, True, INK, None, None, 1.02),
    ])
    textbox(x + 0.22, CARD_Y + 1.00, CARD_W - 0.44, 0.85, [
        (lead, 10, False, INK_SOFT, None, None, 1.24),
    ])
    y = CARD_Y + 1.88
    for b in bullets:
        dot = rect(x + 0.24, y + 0.075, 0.062, 0.062, fill=RGBColor.from_string(accent), radius=0.031)
        textbox(x + 0.40, y - 0.015, CARD_W - 0.62, 0.72, [
            (b, 9.5, False, INK, None, None, 1.2),
        ])
        y += 0.50

# ── Cobertura dos processos de OTC Derivatives ───────────────────────────────
# Comparação pedida pela mesa: quanto do universo de 124 pontos de processo
# cada ferramenta automatiza — o OTC Tracker (target e o já entregue) contra o
# que a tech já desenvolveu (Cockpit, Inoa, Cockpit + AEVO + Registration).
# Tracker nas cores da marca, stack legada em cinza — a barra é a comparação,
# não decoração. Os números vêm dos quadros da mesa (ago/2026).
GRAY_BAR = '9AA3B8'
COBERTURA = [
    ('OTC Tracker — Target',          87.10, A2),
    ('OTC Tracker — Today',           27.42, A1),
    ('Inoa',                          14.52, GRAY_BAR),
    ('Cockpit + AEVO + Registration', 12.10, GRAY_BAR),
    ('Cockpit',                        8.06, GRAY_BAR),
]
COV_Y, COV_H = 5.22, 0.78
COV_LBL_W = 1.75                     # bloco do rótulo à esquerda
COV_IW = (12.09 - COV_LBL_W - 0.25) / len(COBERTURA)
rect(0.62, COV_Y, 12.09, COV_H, fill=CARD_BG, radius=0.17, line=CARD_BD, line_w=Pt(0.75))
textbox(0.92, COV_Y + 0.12, COV_LBL_W - 0.35, 0.42, [
    ('PROCESS COVERAGE', 9, True, INK, None, None, 1.0),
    ('OTC Derivatives', 8, False, INK_SOFT, 2, None, 1.0),
])
for i, (nome, pct, accent) in enumerate(COBERTURA):
    ix = 0.62 + COV_LBL_W + i * COV_IW
    # nome em até duas linhas (com 5 itens a coluna estreitou)
    textbox(ix, COV_Y + 0.08, COV_IW - 0.12, 0.32, [
        (nome, 7.5, False, INK, None, None, 1.05),
    ])
    track_w = COV_IW - 0.70
    rect(ix, COV_Y + 0.42, track_w, 0.085, fill=RGBColor(0xEC, 0xEE, 0xF5), radius=0.042)
    rect(ix, COV_Y + 0.42, max(track_w * pct / 100.0, 0.06), 0.085,
         fill=RGBColor.from_string(accent), radius=0.042)
    textbox(ix + track_w + 0.06, COV_Y + 0.35, 0.58, 0.25, [
        ('%.2f%%' % pct, 8.5, True, RGBColor.from_string(accent), None, None, 1.0),
    ])
textbox(0.92, COV_Y + 0.585, 11.5, 0.18, [
    ('The % is the share of the 124-point OTC Derivatives process universe each tool '
     'automates end-to-end (or nearly so) — everything outside it is still a manual process.',
     7.5, False, INK_FAINT, None, None, 1.0),
])

# ── Rodapé: o que a solução conversa e o que ela vigia ───────────────────────
FOOT_Y, FOOT_H = 6.08, 1.00
foot = rect(0.62, FOOT_Y, 12.09, FOOT_H, fill=CARD_BG, radius=0.17, line=CARD_BD, line_w=Pt(0.75))

# Com QUATRO colunas o texto de cada uma tem de ser mais curto que o das três
# antigas — a largura cai de 4,03" para 3,02" e o corpo não pode passar de
# ~3 linhas, senão vaza do cartão do rodapé.
COLUNAS = [
    ('Connected to', A1,
     'Athena  ·  B3 / CETIP  ·  Reference Data  ·  Electronic Inventory  ·  '
     'SSO + 2FA  ·  E-mail and push'),
    ('Reconciles', A2,
     'FX options: CETIP × Athena end-of-day  ·  Pay / Rec  ·  Comitente  ·  '
     'B3 return files'),
    ('Keeps in sight', A3,
     'New Deals and Confirmations Monitors  ·  KPI and daily metrics  ·  '
     'Alerts when something is late'),
    ('Benefits delivered', A4,
     'Several Intelligent Solutions consolidated in one place  ·  '
     '0.5 FTE reduction already delivered'),
]
CW = 12.09 / len(COLUNAS)
for i, (rotulo, accent, corpo) in enumerate(COLUNAS):
    cx = 0.62 + i * CW
    if i:
        rect(cx, FOOT_Y + 0.22, 0.008, FOOT_H - 0.44, fill=RGBColor(0xEC, 0xEE, 0xF5))
    textbox(cx + 0.30, FOOT_Y + 0.22, CW - 0.55, 0.3, [
        (rotulo.upper(), 9, True, RGBColor.from_string(accent), None, None, 1.0),
    ])
    textbox(cx + 0.30, FOOT_Y + 0.52, CW - 0.55, 0.62, [
        (corpo, 9.5, False, INK_SOFT, None, None, 1.22),
    ])

# ── Assinatura ───────────────────────────────────────────────────────────────
textbox(0.62, 7.18, 12.09, 0.3, [
    ('Internal use — Brazil OTC Operations, J.P. Morgan', 8.5, False, INK_FAINT, None, None, 1.0),
])

# ══ Página 2 — Process Coverage Map ══════════════════════════════════════════
# Os seis quadros da mesa (universo, Cockpit, Inoa, Cockpit+AEVO+Registration,
# Tracker hoje e Tracker target) condensados numa tabela só: cada célula
# processo × produto vira um chip colorido pelo ESTADO — automatizado no
# Tracker hoje, só nas ferramentas legadas, no target do Tracker, ou manual e
# fora do target. É o que deixa a comparação executiva: cinco tabelas de 0/1
# viram um mapa, e a legenda carrega as contagens. Os totais de cada quadro
# são conferidos na geração (asserção abaixo): transcrição de planilha é onde
# um dígito escapa sem ninguém ver.
PROCESSOS = ['Registro', 'Aviso Liquidação', 'Calculo Imposto', 'Confirmação',
             'Recompra', 'Reconciliação Posicao', 'Pagamentos', 'Emissão de CGD',
             'EA ou Aviso Premio D+0', 'Reconciliação Comitente', 'Controle Pay/Rec',
             'Inventário Eletronico', 'Intrag']
FOLHAS = ['EDG', 'CEM', 'Vanilla', 'FWD Start', 'Other Publisher', 'Commodities',
          'EDG', 'CEM', 'Commodities', '', 'EDG', 'CEM', 'NDF', 'DFW', 'FX Swap',
          'Commodities', 'EDG', 'CEM', 'Commodities']
GRUPOS2 = [('Swap', 0, 2), ('NDF', 2, 4), ('Opção', 6, 3), ('Others', 9, 1),
           ('Swap', 10, 2), ('FXCash', 12, 4), ('Opção', 16, 3)]
LADOS = [('ONSHORE', 0, 10, A1), ('OFFSHORE', 10, 9, A3)]

Z9 = [0] * 9


def _mx(*linhas):
    return {p: list(v) for p, v in zip(PROCESSOS, linhas)}


UNIVERSO = _mx(
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + [0, 1, 1, 1, 1, 1, 0, 1, 1],
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + Z9,
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + Z9,
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + [1, 1, 1, 1, 1, 1, 1, 1, 1],
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + [1, 1, 1, 1, 1, 1, 1, 1, 1],
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + [0, 1, 1, 1, 1, 1, 0, 1, 1],
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + Z9,
    [0] * 9 + [1] + Z9,
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + Z9,
    [0] * 9 + [1] + Z9,
    [0] * 9 + [1] + Z9,
    [0] * 9 + [1] + Z9,
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + [0, 1, 1, 1, 1, 1, 0, 1, 1])
TARGET = _mx(
    [0.5, 0.5, 1, 1, 1, 1, 1, 1, 1, 0] + [0, 1, 1, 1, 1, 1, 0, 1, 1],
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + Z9,
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + Z9,
    [0.5, 0.5, 0, 1, 0, 1, 1, 1, 1, 0] + [1, 1, 1, 1, 1, 1, 1, 1, 1],
    [1, 1, 0, 0, 0, 1, 1, 1, 1, 0] + [1, 1, 1, 1, 1, 1, 1, 1, 1],
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + [0, 1, 1, 1, 1, 1, 0, 1, 1],
    [0] * 19,
    [0] * 9 + [1] + Z9,
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + Z9,
    [0] * 9 + [1] + Z9,
    [0] * 9 + [1] + Z9,
    [0] * 9 + [1] + Z9,
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + [0, 1, 1, 1, 1, 1, 0, 1, 1])
HOJE = _mx(
    [0, 0, 1, 1, 1, 1, 0, 1, 1, 0] + Z9,
    [1, 1, 1, 1, 1, 1, 0, 0, 0, 0] + Z9,
    [0, 1, 1, 1, 1, 1, 0, 0, 1, 0] + Z9,
    [0, 0, 0, 1, 0, 1, 0, 1, 1, 0] + Z9,
    [0] * 19,
    [0, 0, 0, 0, 0, 0, 0, 1, 0, 0] + Z9,
    [0] * 19,
    [0] * 19,
    [0, 0, 0, 0, 0, 0, 0, 1, 1, 0] + Z9,
    [0] * 9 + [1] + Z9,
    [0] * 9 + [1] + Z9,
    [0] * 9 + [1] + Z9,
    [0, 0, 1, 1, 1, 1, 0, 1, 1, 0] + Z9)
COCKPIT = _mx(
    [0, 0, 1, 0, 0, 0, 0, 0, 0, 0] + Z9,
    [0, 0, 1, 1, 1, 0, 0, 0, 0, 0] + Z9,
    [0, 0, 1, 1, 1, 0, 0, 0, 0, 0] + Z9,
    [0] * 19, [0] * 19, [0] * 19,
    [0, 0, 1, 1, 1, 0, 0, 0, 0, 0] + Z9,
    [0] * 19, [0] * 19, [0] * 19, [0] * 19, [0] * 19, [0] * 19)
INOA = _mx(
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + Z9,
    [0] * 19, [0] * 19, [0] * 19, [0] * 19,
    [1, 1, 1, 1, 1, 1, 1, 1, 1, 0] + Z9,
    [0] * 19, [0] * 19, [0] * 19, [0] * 19, [0] * 19, [0] * 19, [0] * 19)
CKAR = _mx(
    [0, 0, 1, 0, 0, 0, 0, 0, 0, 0] + Z9,
    [0, 0, 1, 1, 1, 0, 0, 0, 0, 0] + Z9,
    [0, 0, 1, 1, 1, 0, 0, 0, 0, 0] + Z9,
    [0, 0, 1, 0, 1, 0, 0, 0, 0, 0] + Z9,
    [0, 0, 1, 1, 1, 0, 0, 0, 0, 0] + Z9,
    [0] * 19,
    [0, 0, 1, 1, 1, 0, 0, 0, 0, 0] + Z9,
    [0] * 19, [0] * 19, [0] * 19, [0] * 19, [0] * 19, [0] * 19)

# Transcrição conferida contra o Total Geral de cada quadro da mesa.
for _nome, _mxx, _esp in (('universo', UNIVERSO, 124), ('target', TARGET, 108),
                          ('hoje', HOJE, 34), ('cockpit', COCKPIT, 10),
                          ('inoa', INOA, 18), ('ck+aevo+reg', CKAR, 15)):
    _soma = sum(sum(v) for v in _mxx.values())
    if _soma != _esp:
        raise SystemExit('ERRO: quadro {} soma {} (esperado {})'.format(_nome, _soma, _esp))

# Estado de cada célula do universo, na ordem executiva: o que o Tracker já
# faz vence; senão o que as legadas fazem; senão o plano; senão é manual.
C_TODAY, C_LEGACY = '0066CC', 'B9C0D0'
C_TARGET, C_TARGET_HALF = 'CBBCF6', '8B5CF6'
C_MANUAL = 'EFCDD6'


def _estado(p, c):
    if not UNIVERSO[p][c]:
        return None
    if HOJE[p][c]:
        return 'today'
    if COCKPIT[p][c] or INOA[p][c] or CKAR[p][c]:
        return 'legacy'
    if TARGET[p][c]:
        return 'target-half' if TARGET[p][c] < 1 else 'target'
    return 'manual'


N2 = {'today': 0, 'legacy': 0, 'target': 0, 'manual': 0}
for _p in PROCESSOS:
    for _c in range(19):
        _e = _estado(_p, _c)
        if _e:
            N2['target' if _e == 'target-half' else _e] += 1

# ── Primitivas da página (compartilhadas com a prévia HTML) ──────────────────
MAP_Y, MAP_H = 1.62, 4.72
LBL_X, LBL_W = 0.90, 1.72
GX0 = 2.68
COLW = (12.71 - 0.22 - GX0) / 19
GX1 = GX0 + 19 * COLW
ROW_Y0, ROWH, CHIP = MAP_Y + 0.90, 0.285, 0.17

P2 = []


def p2rect(x, y, w, h, hexcolor, radius=None):
    P2.append(('rect', x, y, w, h, hexcolor, radius))


def p2text(x, y, w, txt, size, bold=False, hexcolor='0E112A', align='l', spacing=1.0):
    P2.append(('text', x, y, w, txt, size, bold, hexcolor, align, spacing))


P2.append(('card', 0.62, MAP_Y, 12.09, MAP_H))
for lado, c0, nc, accent in LADOS:
    lx = GX0 + c0 * COLW
    p2text(lx, MAP_Y + 0.10, nc * COLW, lado, 7.5, True, '545A72', 'c')
    p2rect(lx + 0.04, MAP_Y + 0.28, nc * COLW - 0.08, 0.016, accent, 0.008)
for grupo, c0, nc in GRUPOS2:
    gx = GX0 + c0 * COLW
    p2text(gx, MAP_Y + 0.335, nc * COLW, grupo, 6.5, True, '0E112A', 'c')
    p2rect(gx + 0.05, MAP_Y + 0.50, nc * COLW - 0.10, 0.01, 'D9DDE8')
for c, folha in enumerate(FOLHAS):
    if folha:
        p2text(GX0 + c * COLW, MAP_Y + 0.545, COLW, folha, 6, False, '545A72', 'c', 0.95)
for i, proc in enumerate(PROCESSOS):
    ry = ROW_Y0 + i * ROWH
    if i % 2 == 0:
        p2rect(LBL_X - 0.08, ry, GX1 - LBL_X + 0.08, ROWH, 'F3F5FA', 0.05)
    p2text(LBL_X, ry + 0.045, LBL_W, proc, 7.5, False, '0E112A')
    for c in range(19):
        e = _estado(proc, c)
        if not e:
            continue
        cx = GX0 + c * COLW + (COLW - CHIP) / 2.0
        cy = ry + (ROWH - CHIP) / 2.0
        if e == 'today':
            p2rect(cx, cy, CHIP, CHIP, C_TODAY, 0.05)
        elif e == 'legacy':
            p2rect(cx, cy, CHIP, CHIP, C_LEGACY, 0.05)
        elif e == 'target':
            p2rect(cx, cy, CHIP, CHIP, C_TARGET, 0.05)
        elif e == 'target-half':
            # metade forte sobre o chip claro = cobertura parcial (0,5) no target
            p2rect(cx, cy, CHIP, CHIP, C_TARGET, 0.05)
            p2rect(cx, cy, CHIP / 2.0, CHIP, C_TARGET_HALF, 0.05)
        else:
            p2rect(cx, cy, CHIP, CHIP, C_MANUAL, 0.05)

# Legenda com as contagens — é ela que traduz o mapa de volta para os números
# dos quadros (34 + 19 + … = 124).
LEG_Y = MAP_Y + MAP_H + 0.16
LEGENDA2 = [
    (C_TODAY,  None,          'Automated in OTC Tracker today ({})'.format(N2['today'])),
    (C_LEGACY, None,          'Legacy tools — Cockpit / Inoa / AEVO ({})'.format(N2['legacy'])),
    (C_TARGET, None,          'OTC Tracker target ({})'.format(N2['target'])),
    (C_MANUAL, None,          'Manual, outside the target ({})'.format(N2['manual'])),
    (C_TARGET, C_TARGET_HALF, 'half = partial in target'),
]
_lx = 0.62
for cor, meia, rotulo in LEGENDA2:
    p2rect(_lx, LEG_Y + 0.02, 0.15, 0.15, cor, 0.045)
    if meia:
        p2rect(_lx, LEG_Y + 0.02, 0.075, 0.15, meia, 0.045)
    p2text(_lx + 0.22, LEG_Y + 0.015, 3.4, rotulo, 8, False, '545A72')
    _lx += 0.22 + 0.066 * len(rotulo) + 0.16

# ── Desenho da página 2 no PPTX ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])       # os helpers leem o global
rect(0, 0, 13.333, 7.5, fill=PAGE_BG)
band2 = rect(0, 0, 13.333, 1.42)
gradient(band2, [(0, A1), (0.5, A2), (0.8, A3), (1, A4)])
textbox(0.62, 0.24, 8.6, 0.95, [
    ('Process Coverage Map', 30, True, WHITE, None, None, 1.0),
    ('OTC Derivatives — who automates each process, product by product',
     14, False, RGBColor(0xE4, 0xEC, 0xFF), 3, None, 1.0),
])
textbox(7.25, 0.56, 5.45, 0.4, [
    ('Onshore  ·  Offshore   ›   124 process points',
     9.5, True, RGBColor(0xEB, 0xF1, 0xFF), None, PP_ALIGN.RIGHT, 1.0),
])
_P2_ALIGN = {'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER, 'r': PP_ALIGN.RIGHT}
for prim in P2:
    if prim[0] == 'card':
        _, x, y, w, h = prim
        rect(x, y, w, h, fill=CARD_BG, radius=0.17, line=CARD_BD, line_w=Pt(0.75))
    elif prim[0] == 'rect':
        _, x, y, w, h, hexc, radius = prim
        rect(x, y, w, h, fill=RGBColor.from_string(hexc), radius=radius)
    else:
        _, x, y, w, txt, size, bold, hexc, align, spacing = prim
        textbox(x, y, w, 0.4, [
            (txt, size, bold, RGBColor.from_string(hexc), None, _P2_ALIGN[align], spacing),
        ])
textbox(0.62, 7.18, 12.09, 0.3, [
    ('Internal use — Brazil OTC Operations, J.P. Morgan', 8.5, False, INK_FAINT, None, None, 1.0),
])

prs.save(OUT)
print('OK ->', OUT)

# ── Prévia HTML (opcional) ───────────────────────────────────────────────────
# Só roda com um caminho no argv: `python scripts/build_onepager_pptx.py /tmp/p.html`.
# Serve para CONFERIR composição e transbordo de texto sem abrir o PowerPoint —
# a mesma geometria (as mesmas polegadas × 96 px) desenhada em HTML, a partir do
# mesmo `PILARES`/`COLUNAS`, então as duas não divergem. Fica de fora por padrão
# para não largar um .html na raiz do repo.
import sys

if len(sys.argv) < 2:
    raise SystemExit(0)
PREVIEW = sys.argv[1]
PX = 96


def px(v):
    return '%.1fpx' % (v * PX)


bl = []
for i, (accent, titulo, lead, bullets) in enumerate(PILARES):
    x = X0 + i * (CARD_W + GAP)
    itens = ''.join(
        '<div style="position:absolute;left:%s;top:%s;width:%s">'
        '<span style="position:absolute;left:-%s;top:6px;width:6px;height:6px;'
        'border-radius:3px;background:#%s"></span>'
        '<div style="font-size:9.5pt;line-height:1.2;color:#0E112A">%s</div></div>'
        % (px(0.40), px(1.88 - 0.015 + j * 0.50), px(CARD_W - 0.62), px(0.16), accent, b)
        for j, b in enumerate(bullets))
    bl.append(
        '<div style="position:absolute;left:%s;top:%s;width:%s;height:%s;background:#fff;'
        'border:1px solid #E3E6EF;border-radius:%s;box-sizing:border-box">'
        '<div style="position:absolute;left:%s;top:%s;width:%s;height:%s;border-radius:3px;background:#%s"></div>'
        '<div style="position:absolute;left:%s;top:%s;width:%s;font:700 15pt %s;color:#0E112A;line-height:1.02">%s</div>'
        '<div style="position:absolute;left:%s;top:%s;width:%s;font-size:10pt;color:#545A72;line-height:1.24">%s</div>'
        '%s</div>'
        % (px(x), px(CARD_Y), px(CARD_W), px(CARD_H), px(0.17),
           px(0.22), px(0.28), px(0.52), px(0.075), accent,
           px(0.22), px(0.48), px(CARD_W - 0.44), FONT, titulo,
           px(0.22), px(1.00), px(CARD_W - 0.44), lead, itens))

# faixa de cobertura — mesma geometria e mesmo `COBERTURA` do slide
cov_items = []
for i, (nome, pct, accent) in enumerate(COBERTURA):
    ix = COV_LBL_W + i * COV_IW
    track_w = COV_IW - 0.70
    cov_items.append(
        '<div style="position:absolute;left:%s;top:0;width:%s">'
        '<div style="position:absolute;left:0;top:%s;width:%s;font-size:7.5pt;color:#0E112A;line-height:1.05">%s</div>'
        '<div style="position:absolute;left:0;top:%s;width:%s;height:%s;border-radius:4px;background:#ECEEF5"></div>'
        '<div style="position:absolute;left:0;top:%s;width:%s;height:%s;border-radius:4px;background:#%s"></div>'
        '<div style="position:absolute;left:%s;top:%s;font:700 8.5pt %s;color:#%s">%s</div>'
        '</div>'
        % (px(ix), px(COV_IW),
           px(0.08), px(COV_IW - 0.12), nome,
           px(0.42), px(track_w), px(0.085),
           px(0.42), px(max(track_w * pct / 100.0, 0.06)), px(0.085), accent,
           px(track_w + 0.06), px(0.35), FONT, accent, '%.2f%%' % pct))
cov = (
    '<div style="position:absolute;left:%s;top:%s;width:%s;height:%s;background:#fff;'
    'border:1px solid #E3E6EF;border-radius:%s;box-sizing:border-box">'
    '<div style="position:absolute;left:%s;top:%s;font:700 9pt %s;color:#0E112A;letter-spacing:.04em">PROCESS COVERAGE</div>'
    '<div style="position:absolute;left:%s;top:%s;font-size:8pt;color:#545A72">OTC Derivatives</div>'
    '<div style="position:absolute;left:%s;top:%s;width:%s;font-size:7.5pt;color:#7B8299">'
    'The %% is the share of the 124-point OTC Derivatives process universe each tool automates '
    'end-to-end (or nearly so) — everything outside it is still a manual process.</div>'
    '%s</div>'
    % (px(0.62), px(COV_Y), px(12.09), px(COV_H), px(0.17),
       px(0.30), px(0.12), FONT,
       px(0.30), px(0.32),
       px(0.30), px(0.585), px(11.5),
       ''.join(cov_items)))

cols = []
for i, (rotulo, accent, corpo) in enumerate(COLUNAS):
    cx = i * CW
    cols.append(
        '<div style="position:absolute;left:%s;top:%s;width:%s">'
        '%s'
        '<div style="position:absolute;left:%s;top:%s;font:700 9pt %s;color:#%s;letter-spacing:.04em">%s</div>'
        '<div style="position:absolute;left:%s;top:%s;width:%s;font-size:9.5pt;color:#545A72;line-height:1.22">%s</div>'
        '</div>'
        % (px(cx), '0px', px(CW),
           ('' if not i else '<div style="position:absolute;left:0;top:%s;width:1px;height:%s;background:#ECEEF5"></div>'
            % (px(0.22), px(FOOT_H - 0.44))),
           px(0.30), px(0.22), FONT, accent, rotulo.upper(),
           px(0.30), px(0.52), px(CW - 0.55), corpo))

html = (
    '<div style="position:relative;width:%s;height:%s;background:#F6F7FB;'
    'font-family:%s,Helvetica,Arial,sans-serif;overflow:hidden">'
    '<div style="position:absolute;inset:0 0 auto 0;height:%s;'
    'background:linear-gradient(100deg,#0066CC,#5E5CE6 50%%,#8B5CF6 80%%,#D946EF)"></div>'
    '<div style="position:absolute;left:%s;top:%s;width:%s;font:700 30pt %s;color:#fff;line-height:1">OTC Tracker</div>'
    '<div style="position:absolute;left:%s;top:%s;width:%s;font-size:14pt;color:#E4ECFF">'
    'One platform for the full OTC derivatives lifecycle — Brazil OTC Operations</div>'
    '<div style="position:absolute;left:%s;top:%s;width:%s;text-align:right;font:700 9.5pt %s;'
    'color:#EBF1FF;line-height:1.18">Trade capture &nbsp;›&nbsp; B3 / CETIP registration &nbsp;›&nbsp; '
    'Client confirmation &nbsp;›&nbsp; Settlement</div>'
    '%s'
    '%s'
    '<div style="position:absolute;left:%s;top:%s;width:%s;height:%s;background:#fff;border:1px solid #E3E6EF;'
    'border-radius:%s;box-sizing:border-box">%s</div>'
    '<div style="position:absolute;left:%s;top:%s;font-size:8.5pt;color:#7B8299">'
    'Internal use — Brazil OTC Operations, J.P. Morgan</div>'
    '</div>'
    % (px(13.333), px(7.5), FONT, px(1.42),
       px(0.62), px(0.24), px(8.6), FONT,
       px(0.62), px(0.80), px(8.6),
       px(7.25), px(0.56), px(5.45), FONT,
       ''.join(bl), cov,
       px(0.62), px(FOOT_Y), px(12.09), px(FOOT_H), px(0.17), ''.join(cols),
       px(0.62), px(7.18)))

# página 2 — desenhada das MESMAS primitivas P2 do slide, então não divergem
p2divs = []
for prim in P2:
    if prim[0] == 'card':
        _, x, y, w, h = prim
        p2divs.append('<div style="position:absolute;left:%s;top:%s;width:%s;height:%s;'
                      'background:#fff;border:1px solid #E3E6EF;border-radius:%s;'
                      'box-sizing:border-box"></div>' % (px(x), px(y), px(w), px(h), px(0.17)))
    elif prim[0] == 'rect':
        _, x, y, w, h, hexc, radius = prim
        p2divs.append('<div style="position:absolute;left:%s;top:%s;width:%s;height:%s;'
                      'background:#%s;border-radius:%s"></div>'
                      % (px(x), px(y), px(w), px(h), hexc, px(radius or 0)))
    else:
        _, x, y, w, txt, size, bold, hexc, align, spacing = prim
        p2divs.append('<div style="position:absolute;left:%s;top:%s;width:%s;'
                      'font:%s %spt %s;color:#%s;text-align:%s;line-height:%s">%s</div>'
                      % (px(x), px(y), px(w), '700' if bold else '400', size, FONT, hexc,
                         {'l': 'left', 'c': 'center', 'r': 'right'}[align], spacing, txt))
html2 = (
    '<div style="position:relative;width:%s;height:%s;background:#F6F7FB;'
    'font-family:%s,Helvetica,Arial,sans-serif;overflow:hidden">'
    '<div style="position:absolute;inset:0 0 auto 0;height:%s;'
    'background:linear-gradient(100deg,#0066CC,#5E5CE6 50%%,#8B5CF6 80%%,#D946EF)"></div>'
    '<div style="position:absolute;left:%s;top:%s;width:%s;font:700 30pt %s;color:#fff;line-height:1">Process Coverage Map</div>'
    '<div style="position:absolute;left:%s;top:%s;width:%s;font-size:14pt;color:#E4ECFF">'
    'OTC Derivatives — who automates each process, product by product</div>'
    '<div style="position:absolute;left:%s;top:%s;width:%s;text-align:right;font:700 9.5pt %s;'
    'color:#EBF1FF">Onshore &nbsp;·&nbsp; Offshore &nbsp;&nbsp;›&nbsp;&nbsp; 124 process points</div>'
    '%s'
    '<div style="position:absolute;left:%s;top:%s;font-size:8.5pt;color:#7B8299">'
    'Internal use — Brazil OTC Operations, J.P. Morgan</div>'
    '</div>'
    % (px(13.333), px(7.5), FONT, px(1.42),
       px(0.62), px(0.24), px(9.6), FONT,
       px(0.62), px(0.80), px(9.6),
       px(7.25), px(0.56), px(5.45), FONT,
       ''.join(p2divs),
       px(0.62), px(7.18)))

with open(PREVIEW, 'w', encoding='utf-8') as fh:
    fh.write('<!doctype html><meta charset="utf-8"><body style="margin:0">'
             + html + '<div style="height:24px"></div>' + html2)
print('prévia ->', PREVIEW)
